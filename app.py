import streamlit as st
import json
import os
import io
import csv # CSV処理ライブラリ
import time # ファイルアップロード後の待機用
from pydantic import BaseModel, Field, ValidationError, ConfigDict # ConfigDictをインポート
from typing import Optional, Literal, Dict, Any, List, Union # Unionを追加
import re # 正規表現モジュールを追加

# 外部ライブラリ
import pypdf # PDF処理ライブラリ
import docx # DOCX処理ライブラリ (python-docx)
import openpyxl # XLSX処理ライブラリ
from pptx import Presentation # PPTX処理ライブラリ (python-pptx)

# ----------------------------------------------------------------------
# ⚠️ 注意: Gemini API依存を削除するため、Google Geminiのライブラリは使用しません。
# ----------------------------------------------------------------------

# ----------------------------------------------------------------------
# 1. データ構造の定義 (Pydanticで継続)
# ----------------------------------------------------------------------

# 論文データ => 著者付き文書データに名称変更
class AuthorData(BaseModel):
    # year: str = Field(description="出版年西暦 (例: 2024)") # 年号は必須ではないためロジックでのみ利用
    author: str = Field(description="主要著者名。カンマ区切りで記述してください。")
    title: str = Field(description="文書のタイトル。")

# 請求書・領収書データ
class InvoiceData(BaseModel):
    invoice_date: str = Field(description="発行日。YYYY-MM-DD形式に変換してください。")
    invoice_amount: str = Field(description="合計金額。数字と通貨記号を含んだ元の文字列。")
    invoice_issuer: str = Field(description="発行元/発行者名。")
    invoice_subject: str = Field(description="請求書/領収書の件名。")

# その他データ
class OtherData(BaseModel):
    title: str = Field(description="ファイル内容を最もよく表す、AIが推測したタイトル。")

# AIコアからの最終応答スキーマ
Category = Literal["論文", "請求書・領収書", "その他", "不明"]

class AICoreResponse(BaseModel):
    # 余分な入力を無視する設定
    model_config = ConfigDict(extra='ignore')

    category: Category = Field(description="ファイルの分類カテゴリ。必須。取りうる値: 論文, 請求書・領収書, その他, 不明")
    # AuthorData (旧 PaperData) を使用
    extracted_data: Optional[Union[AuthorData, InvoiceData, OtherData, Dict[str, Any]]] = Field( 
        None, 
        description="分類に応じた抽出データを含むオブジェクト。不明の場合は null にしてください。"
    )
    reasoning: str = Field(description="その分類と抽出を行った根拠。")
    transcript: Optional[str] = Field(None, description="音声ファイルが入力された場合の文字起こし結果。")

# ----------------------------------------------------------------------
# 2. バックエンド処理機能 (ファイル抽出とローカルAI連携)
# ----------------------------------------------------------------------

# --- (extract_text 関数は変更なし) ---
def extract_text(uploaded_file: st.runtime.uploaded_file_manager.UploadedFile) -> tuple[str, bool]:
    """
    ファイル形式に応じてテキストを抽出する関数。
    音声ファイルは「文字起こしが必要」としてフラグ (is_asr=True) を返す。
    """
    file_ext = uploaded_file.name.split('.')[-1].lower()
    
    # 対応ファイル形式のチェック
    supported_extensions = ['pdf', 'docx', 'xlsx', 'pptx', 'csv', 'mp3', 'wav', 'm4a']
    if file_ext not in supported_extensions:
        return f"ファイル形式 '{file_ext}' は対応していません。", False

    # --- 音声ファイル処理 (フラグを返す) ---
    if file_ext in ['mp3', 'wav', 'm4a']:
        st.info(f"🔊 音声ファイル ({uploaded_file.name}): ローカルASR処理モックを使用します。")
        # ⚠️ Geminiを使わないため、ASRはローカルでモックとして処理する
        return uploaded_file.name, True 

    # --- PDF 処理 (安定性強化) ---
    if file_ext == 'pdf':
        try:
            st.info(f"📄 PDFファイル ({uploaded_file.name}): テキスト抽出を実行中...")
            pdf_reader = pypdf.PdfReader(uploaded_file)
            text_content = ""
            for page in pdf_reader.pages:
                try:
                    text_content += page.extract_text() or ""
                except (TypeError, ValueError) as e:
                    st.warning(f"⚠️ ページ抽出エラー: {e}")
                    continue
                
            if not text_content.strip():
                st.warning("⚠️ PDFからテキストが抽出できませんでした。スキャン画像と見なしてモックOCRテキストを使用します。")
                text_content = "OCR結果: このファイルは2024年4月1日に発行された領収書であり、金額は25,000円です。発行元はABCコンサルティングです。"
            
            return text_content, False
        
        except Exception as e:
            st.error(f"🚨 PDF処理エラー: {e}")
            return f"PDF処理中にエラーが発生しました: {e}", False

    # --- DOCX 処理 ---
    elif file_ext == 'docx':
        try:
            st.info(f"📄 DOCXファイル ({uploaded_file.name}): テキスト抽出を実行中...")
            document = docx.Document(io.BytesIO(uploaded_file.getvalue()))
            text_content = ""
            for paragraph in document.paragraphs:
                text_content += paragraph.text + '\n' 
                
            if not text_content.strip():
                st.warning("⚠️ DOCXからテキストが抽出できませんでした。ファイル内容が空か、読み取りに失敗しました。")
            
            return text_content, False

        except Exception as e:
            st.error(f"🚨 DOCX処理エラー: {e}")
            return f"DOCX処理中にエラーが発生しました: {e}", False

    # --- XLSX 処理 ---
    elif file_ext == 'xlsx':
        try:
            st.info(f"📊 XLSXファイル ({uploaded_file.name}): テキスト抽出を実行中...")
            workbook = openpyxl.load_workbook(uploaded_file, read_only=True)
            text_content = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"\n--- シート: {sheet_name} ---\n"
                
                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                         if cell.value is not None:
                            row_data.append(str(cell.value))
                    if row_data:
                        text_content += ', '.join(row_data) + '\n'
            
            if not text_content.strip():
                st.warning("⚠️ XLSXからテキストが抽出できませんでした。")
            
            return text_content, False

        except Exception as e:
            st.error(f"🚨 XLSX処理エラー: {e}")
            return f"XLSX処理中にエラーが発生しました: {e}", False

    # --- PPTX 処理 (安定性強化) ---
    elif file_ext == 'pptx':
        try:
            st.info(f"🖼️ PPTXファイル ({uploaded_file.name}): テキスト抽出を実行中...")
            presentation = Presentation(uploaded_file)
            text_content = ""
            
            for i, slide in enumerate(presentation.slides):
                text_content += f"\n--- スライド {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame: # テキストフレームの存在チェック
                        text_content += shape.text + '\n'
                    elif shape.has_table:
                        # テーブルセル内のテキストをより確実に取得
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                if cell.text_frame:
                                    row_data.append(cell.text)
                            text_content += ' | '.join(row_data) + '\n'
                    elif shape.has_text_frame: # has_text_frameはtext_frameの有無をチェック
                        text_content += shape.text_frame.text + '\n'

            if not text_content.strip():
                st.warning("⚠️ PPTXからテキストが抽出できませんでした。")

            return text_content, False
        
        except Exception as e:
            st.error(f"🚨 PPTX処理エラー: {e}")
            return f"PPTX処理中にエラーが発生しました: {e}", False

    # --- CSV 処理 ---
    elif file_ext == 'csv':
        try:
            st.info(f"📋 CSVファイル ({uploaded_file.name}): テキスト抽出を実行中...")
            text_stream = io.StringIO(uploaded_file.getvalue().decode('utf-8'))
            reader = csv.reader(text_stream)
            
            text_content = ""
            for row in reader:
                text_content += ', '.join(row) + '\n'

            if not text_content.strip():
                st.warning("⚠️ CSVファイルが空か、読み取りに失敗しました。")

            return text_content, False

        except Exception as e:
            st.error(f"🚨 CSV処理エラー: {e}")
            return f"CSV処理中にエラーが発生しました: {e}", False
            

# --- Gemini API連携を置き換えるローカル分析関数に修正 ---
def analyze_file_content(text_content: str, uploaded_file: st.runtime.uploaded_file_manager.UploadedFile, is_asr: bool) -> AICoreResponse:
    """
    Gemini APIの代わりに、ローカルのルールベースでファイル内容を分析し、構造化データを生成する。
    """
    
    # ------------------------------------------------------------------
    # 0. 音声ファイル処理 (文字起こしと分析)
    # ------------------------------------------------------------------
    if is_asr:
        # ⚠️ 音声ファイル処理のモック (要件 4. ASR)
        transcript = "モック文字起こし: 2023年10月5日、田中商事から15000円の請求書を受領しました。件名はソフトウェアライセンスです。"
        st.info("🔎 **分析開始**: 音声ファイルのため、文字起こし結果（モック）に基づき文書分類を行います。")
        
        # モックの抽出データ（文字起こし結果に基づくと仮定）
        data = InvoiceData(
            invoice_date="2023-10-05",
            invoice_amount="15000円",
            invoice_issuer="田中商事",
            invoice_subject="ソフトウェアライセンス"
        )
        return AICoreResponse(
            category="請求書・領収書",
            extracted_data=data,
            reasoning="音声ファイルが検出されました。ローカルASRモックにより文字起こしを行い、その結果から請求情報（日付、金額、発行元）を検出しました。",
            transcript=transcript
        )

    # 文書ファイルの内容分析 (ルールベース)
    lower_text = text_content.lower()
    first_10_lines = '\n'.join(text_content.split('\n')[:10]).strip() # 先頭10行を分析
    
    # 処理状況の表示
    st.info("🔎 **分析開始**: 文書ファイルの内容をローカルルールでスコアリングします。")
    
    # スコアリング基準
    score_invoice = 0
    score_author_doc = 0 # 論文/著者付き文書のスコア
    
    # ------------------------------------------------------------------
    # 1. 請求書/領収書 ルール (スコアベース)
    # ------------------------------------------------------------------
    
    invoice_keywords = ["請求書", "領収書", "明細", "invoice", "receipt", "合計金額", "御中"]
    if any(keyword in lower_text for keyword in invoice_keywords):
        score_invoice += 5
        st.info(f"→ 請求書キーワード検出 ({score_invoice}点)")
    
    date_match = re.search(r"(\d{4}[-/年]\d{1,2}[-/月]\d{1,2}日?)", first_10_lines)
    amount_match = re.search(r"([¥￥$€£]\s*[\d,]+\.?\d*|[\d,]+\s*(円|yen))", first_10_lines)
    
    if date_match:
        score_invoice += 5 # 日付検出
        st.info(f"→ ヘッダーで日付パターン検出 (+5点, 現在{score_invoice}点)")
    if amount_match:
        score_invoice += 5 # 金額検出
        st.info(f"→ ヘッダーで金額パターン検出 (+5点, 現在{score_invoice}点)")
    
    # ------------------------------------------------------------------
    # 2. 著者付き文書 ルール (スコアベース)
    # ------------------------------------------------------------------
    
    author_doc_keywords = [
        "abstract", "introduction", "author", "year of publication", # 論文キーワード
        "抄録", "緒言", "序論", "著者", "発表年", "研究報告", "キーワード", # 論文キーワード
        "レポート", "Report", "技術資料", "作成者", "執筆者" # 一般的な著者付き文書キーワードを追加
    ]
    if any(keyword in lower_text for keyword in author_doc_keywords):
        score_author_doc += 5
        st.info(f"→ 著者付き文書キーワード検出 ({score_author_doc}点)")
    
    # [修正ポイント] 著者名検出の正規表現を日本語名と英語名に対応させる
    # 修正前: author_pattern = re.search(r"(?:Author|著者|作成者|執筆者)\s*[:]?\s*([A-Z][a-z]+(?:\s*[A-Z][a-z]+)?)\s*(?:\((.+?)\))?", first_10_lines)
    # 修正後: 氏名パターンを「漢字/ひらがな/カタカナ」または「英語名」のいずれかに広く対応させる
    
    # 複雑なパターン: 「著者名（所属）」または「著者名（改行）所属」を捉える。
    # [^:]*?: コロン以外の任意の文字（氏名/タイトル）
    # (?:Author|著者|作成者|執筆者): キーワードのいずれか
    # ([A-Z][a-z]+(?:\s*[A-Z][a-z]+)?): 英語名パターン
    # (?:[\u3005\u3006\u303b\u4e00-\u9faf\u3040-\u309f\u30a0-\u30ff]+) : 日本語名パターン
    
    # 今回は、シンプルに「著者名の後に所属機関名が続くパターン」を幅広くカバーする
    author_pattern_match = re.search(
        r"(?:Author|著者|作成者|執筆者)[:\s]*\s*([^,\n]+?)\s*(\([^\n]+\)|[^\n]*\s*[大|会|学|社|科|院|部|校][^\n]*)", 
        first_10_lines, 
        re.IGNORECASE | re.DOTALL
    )
    
    # 日本語/英語両対応の著者名（氏名のみ）をキャプチャするパターン
    name_pattern = r"([^,\n]+?)" # 氏名は、改行やカンマまで
    
    # 「氏名 + 所属」または「氏名 + 役職」がヘッダーにあるか
    author_pattern = re.search(
        r"(?:Author|著者|作成者|執筆者)[\s:]*?(" + name_pattern + r")\s*(\([^\)]+\)|[^\n]+\s*[大|学|社|会|科|院|部|校][^\n]*?)", 
        first_10_lines, 
        re.IGNORECASE | re.DOTALL
    )
    
    # シンプルな「氏名」単体のパターン (例: 町田佳世子 札幌市立大学デザイン学部)
    # 氏名のパターンを寛容にする: 任意の文字 (\w) を含む、改行やカンマを含まない文字列
    author_simple_pattern = re.search(
        r"(?:Author|著者|作成者|執筆者)[\s:]*?([^\n,]+)", 
        first_10_lines, 
        re.IGNORECASE
    )
    
    # 今回のPDFの形式 ('町田佳世子\n 札幌市立大学デザイン学部')に対応するため、
    # 著者の後に所属機関名（日本語の組織名を含む）が続くパターンを優先する
    
    # 氏名（日本語または英語）をキャプチャするパターン
    name_capture_group = r"([^,\n\s]+(?:\s[^,\n\s]+)*?)"
    
    # 氏名が検出され、その後に所属機関っぽいキーワードが続くパターン
    author_doc_match = re.search(
        r"(?:Author|著者|作成者|執筆者)[\s:]*?" + name_capture_group + r"\s*([^\n]*?大学|[^\n]*?研究室|[^\n]*?株式会社)",
        first_10_lines,
        re.IGNORECASE | re.DOTALL
    )
    
    # 最終的な著者情報検出に使用する変数
    detected_author = None
    if author_doc_match:
        detected_author = author_doc_match.group(1).strip()
        st.info(f"→ **構造的著者情報（{detected_author}）**検出 (+10点, 現在{score_author_doc}点)")
        score_author_doc += 10 # 構造的な著者情報検出
    
    
    year_match = re.search(r"(\d{4})", first_10_lines)
    
    if detected_author: # 著者名が検出された場合
        score_author_doc += 10 # 構造的な著者情報検出（再加算ではなく、確実に10点以上にするための補強）
    
    if year_match and score_author_doc > 0:
        score_author_doc += 3 # 年号が検出され、かつ著者付き文書の可能性が高い場合
        st.info(f"→ ヘッダーで年号パターン検出 (+3点, 現在{score_author_doc}点)")
        
    # ------------------------------------------------------------------
    # 3. 最終判定ロジック
    # ------------------------------------------------------------------
    
    reasoning_detail = f"（著者文書スコア: {score_author_doc}, 請求書スコア: {score_invoice}）"
    
    # 論文/著者付き文書と判定
    if score_author_doc >= 10 and score_author_doc > score_invoice:
        st.success(f"✅ **最終判定**: 著者付き文書（論文/レポート）と決定しました。")
        
        # 抽出ロジック（著者付き文書）
        author = detected_author if detected_author else "著者名不明"
        
        # タイトルはテキストの最初の非空白行とする (最も確実)
        # ただし、最初の行が著者名でないことを確認する必要がある (今回は最初の行がタイトルなのでOKとする)
        title_lines = [line for line in text_content.split('\n') if line.strip()]
        
        # 最初の3行から最も長い行をタイトルと見なすロジック（日本語文書対応）
        title_extracted = os.path.splitext(uploaded_file.name)[0] # 初期値はファイル名
        if len(title_lines) > 0:
            # 最初の数行の最も長いものをタイトルとする
            top_lines = title_lines[:4]
            # 著作権表記（Copyrightなど）やジャーナル名は除外したいが、ここでは最も長いものを採用
            title_extracted = max(top_lines, key=len)
        
        data = AuthorData( # AuthorDataを使用
            author=author,
            title=title_extracted 
        )
        # Yearはリネーム形式から削除したため、抽出データには含めない
        return AICoreResponse(
            category="論文", # 要件定義書の分類カテゴリは「論文」を維持
            extracted_data=data,
            reasoning=f"高度なパターンマッチングにより、著者情報（氏名パターン）とキーワードを検出（{score_author_doc}点）。著者付き文書と判定しました。",
        )

    # 請求書と判定
    elif score_invoice >= 10 and score_invoice >= score_author_doc:
        st.success(f"✅ **最終判定**: 請求書/領収書と決定しました。")

        # 抽出ロジック（請求書）
        invoice_date_raw = date_match.group(1) if date_match else "YYYYMMDD"
        invoice_date = invoice_date_raw.replace('年', '-').replace('月', '-').replace('日', '')
        
        amount_extracted = amount_match.group(0) if amount_match else "0"
        
        data = InvoiceData(
            invoice_date=invoice_date,
            invoice_amount=amount_extracted,
            invoice_issuer="不明な発行元", 
            invoice_subject=uploaded_file.name
        )
        return AICoreResponse(
            category="請求書・領収書",
            extracted_data=data,
            reasoning=f"高度なパターンマッチングにより、請求キーワード、日付、金額（{score_invoice}点）を検出し、請求書と判定しました。{reasoning_detail}",
        )

    # 4. その他/不明
    if text_content.strip():
        st.warning("⚠️ **最終判定**: 特定の文書パターンに一致しませんでした。")
        # テキストがあれば「その他」としてファイル名をタイトルとして提案
        data = OtherData(
            title=os.path.splitext(uploaded_file.name)[0]
        )
        return AICoreResponse(
            category="その他",
            extracted_data=data,
            reasoning=f"特定の文書パターン（著者文書、請求書）に一致しませんでした。{reasoning_detail} ファイル名を元にリネームします。"
        )
    else:
        st.error("❌ **最終判定**: ファイル内容が空です。")
        # テキストが空の場合
        return AICoreResponse(
            category="不明",
            extracted_data=None,
            reasoning="ファイルから内容（テキスト）を抽出できませんでした。"
        )


def apply_rename_rule(ai_response: AICoreResponse, original_name: str) -> str:
    """
    要件 6 に基づき、AIの応答からリネーム後のファイル名を生成する。
    """
    base_name, ext = os.path.splitext(original_name)
    category = ai_response.category
    
    # データを dict 形式で取得。extracted_data が None の場合は空の dict を使用
    # モック処理なので、Pydanticモデルから直接 dict に変換 (エラー回避のため)
    data = ai_response.extracted_data.model_dump() if ai_response.extracted_data else {} 

    # ファイル名に使用できない文字を削除/置換するヘルパー関数
    def sanitize_filename(name: str) -> str:
        safe_name = name.replace(' ', '_')
        # ファイル名に使用可能な文字のみを許可
        return ''.join(c for c in safe_name if c.isalnum() or c in '._-')

    # 4. 不明: リネームスキップ
    if category == "不明":
        st.warning("⚠️ カテゴリが「不明」のため、リネーム処理はスキップされました。")
        return original_name

    # 1. 論文 (要件 6.1) -> 著者付き文書としてリネーム (年号なし)
    elif category == "論文":
        # year = data.get("year", "YYYY") # 年号は使用しない
        authors = data.get("author", "著者名不明")
        title = data.get("title", "タイトル不明")

        authors_short = authors[:15] if len(authors) > 15 else authors
        # 最大50字の制限は著者名とタイトルで適用
        max_total_len = 50 - 1 # 1は区切り文字 '_' の数
        
        # 著者を15字に制限後、残りの文字をタイトルに割り当てる
        max_title_len = max_total_len - len(authors_short)
        title_short = title[:max(0, max_title_len)]

        # 命名規則: 著者名_タイトル
        new_name_raw = f"{authors_short}_{title_short}".strip('_')
        return f"{sanitize_filename(new_name_raw)}{ext}"

    # 2. 請求書・領収書 (要件 6.2)
    elif category == "請求書・領収書":
        date_str_raw = data.get("invoice_date", "YYYYMMDD")
        # 日付に含まれる数字だけを取り出し、8桁に制限
        date_str = ''.join(filter(str.isdigit, date_str_raw))[:8]

        issuer = data.get("invoice_issuer", "発行元不明")[:15]
        
        amount_raw = data.get("invoice_amount", "0")
        # 金額に含まれる数字だけを取り出し、カンマや通貨記号を削除
        amount = ''.join(filter(str.isdigit, amount_raw)) or "0" 
        
        subject = data.get("invoice_subject", "件名なし")[:15]

        new_name_raw = f"{date_str}_{issuer}_{amount}_{subject}"
        return f"{sanitize_filename(new_name_raw)}{ext}"

    # 3. その他 (要件 6.3)
    elif category == "その他":
        title = data.get("title", "AI推測タイトル")[:30]
        return f"{sanitize_filename(title)}{ext}"
    
    # 予期せぬ分類エラー
    else:
        st.error(f"🚨 リネームルール適用エラー: カテゴリ '{category}' またはデータ構造が不正です。元のファイル名を返します。")
        return original_name

# ----------------------------------------------------------------------
# 3. Streamlit UI定義 (要件 3)
# ----------------------------------------------------------------------

# ページ設定
st.set_page_config(page_title="🤖 AIスマートファイルリネームシステム (Local Mode)", layout="wide")

## サイドバー
with st.sidebar:
    st.header("⚙️ システム設定")
    st.markdown("""
    **動作モード:** 誰でも使えるローカルルールベースモード
    
    *Gemini APIを使用しないため、APIキーは不要です。*
    *文書分析にはPythonの正規表現を使用します。*
    *音声文字起こしは固定のモック応答となります。*
    """)
    
    st.markdown("---")
    st.subheader("対応ファイル形式 (要件 4)")
    st.markdown("""
    * **文書**: PDF, DOCX, XLSX, PPTX, CSV
    * **音声**: MP3, WAV, M4A (モック)
    """)

## メインエリア
st.title("🤖 AIスマートファイルリネームシステム (Local Mode)")
st.caption("アップロードされたファイルの内容をローカルのルールで分析し、自動リネームを行います。")

# ファイルアップロードエリア (要件 3)
uploaded_files = st.file_uploader(
    "ファイルをアップロード (複数選択可)", 
    type=['pdf', 'docx', 'xlsx', 'pptx', 'csv', 'mp3', 'wav', 'm4a'],
    accept_multiple_files=True
)

if uploaded_files:
    if st.button("🚀 ローカルリネーム・文字起こしを実行", use_container_width=True):
        
        # 処理状況の表示 (要件 3)
        st.subheader("📊 処理結果")
        results: List[Dict[str, Any]] = []
        
        progress_bar = st.progress(0)
        
        with st.empty(): 
            for i, uploaded_file in enumerate(uploaded_files):
                
                progress_bar.progress((i + 1) / len(uploaded_files))
                st.info(f"👉 **{uploaded_file.name}** の処理を開始...")
                
                # 1. テキスト抽出/ASR判定
                text_content, is_asr = extract_text(uploaded_file)
                
                if "対応していません" in text_content or "エラー" in text_content:
                    results.append({
                        "オリジナルファイル名": uploaded_file.name,
                        "処理状況": "スキップ/エラー",
                        "分類カテゴリ": "-",
                        "リネーム後ファイル名": uploaded_file.name,
                    })
                    continue
                
                # 2. ローカルAIコア連携
                ai_response = analyze_file_content(text_content, uploaded_file, is_asr)
                
                if ai_response.category == "不明":
                    st.error(f"❌ ファイル {uploaded_file.name} の処理に失敗しました。理由: {ai_response.reasoning}")

                # 3. リネームルール適用 (要件 6)
                new_filename = apply_rename_rule(ai_response, uploaded_file.name)
                
                # 4. 結果の記録とダウンロードボタンの設置
                result_data = {
                    "オリジナルファイル名": uploaded_file.name,
                    "処理状況": "完了" if ai_response.category != "不明" else "失敗",
                    "分類カテゴリ": ai_response.category,
                    "リネーム後ファイル名": new_filename,
                }
                results.append(result_data)
                
                st.markdown(f"**結果 ({uploaded_file.name})**:")
                
                col1, col2, col3 = st.columns([1, 1, 2])
                
                with col1:
                    st.download_button(
                        label=f"💾 {new_filename} をダウンロード",
                        data=uploaded_file.getvalue(), 
                        file_name=new_filename,
                        mime=uploaded_file.type,
                        key=f"download_renamed_{uploaded_file.name}"
                    )

                if is_asr and ai_response.transcript:
                    with col2:
                        asr_file_name = f"{os.path.splitext(uploaded_file.name)[0]}.txt"
                        st.download_button(
                            label=f"📝 {asr_file_name} ダウンロード",
                            data=ai_response.transcript,
                            file_name=asr_file_name,
                            mime="text/plain",
                            key=f"download_asr_{uploaded_file.name}"
                        )
                
                with col3:
                    st.caption(f"分類: **{ai_response.category}** | 理由: {ai_response.reasoning}")

            st.success("✅ 全ファイルの処理が完了しました！")

        st.dataframe(results, use_container_width=True)
        
        st.markdown("---")
        st.subheader("💡 最終分析結果 (構造化データ)")
        if 'ai_response' in locals() and ai_response:
            # Pydanticモデルを辞書に変換して表示
            st.json(ai_response.model_dump())
        else:
            st.write("ファイルが処理されていません。")
